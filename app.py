import os
import io
from flask import Flask, render_template_string, request, session
from flask_sqlalchemy import SQLAlchemy
#from pyngrok import ngrok
from transformers import pipeline, BartTokenizer, BartForConditionalGeneration
from PIL import Image
import fitz  # PyMuPDF
import pdfplumber
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import torch
import torchvision.transforms as transforms
from torchvision import models
import pytesseract
import groq
from spire.doc import *
from spire.doc.common import *
# import docx2txt
# !pip install python-docx
import re
# Import required libraries
from docx import Document
from google.colab import files
from datetime import datetime
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize

# Step 3: Define the Flask app
app = Flask(__name__)
app.secret_key = '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'  # Required for session management

# SQLAlchemy configuration
app.config['SECRET_KEY'] = 'your_secret_key'
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///summaries.db'
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
db = SQLAlchemy(app)

# Define the Summary model
class Summary(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    content = db.Column(db.Text, nullable=False)

class QuestionAnswer(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    question = db.Column(db.String(500), nullable=False)
    answer = db.Column(db.Text, nullable=False)
    timestamp = db.Column(db.DateTime, default=datetime.utcnow)

    def __repr__(self):
        return f'<QuestionAnswer {self.id}>'

@app.before_first_request
def create_tables():
    db.create_all()

# Ensure the uploads directory exists
os.makedirs("uploads", exist_ok=True)

# Define the labels for graph types
graph_labels = ["bar", "line", "pie", "scatter"]

# Load and configure ResNet50 model for graph classification
graph_model = models.resnet50(pretrained=True)
num_classes = len(graph_labels)
graph_model.fc = torch.nn.Linear(graph_model.fc.in_features, num_classes)
graph_model.eval()

def preprocess_text(text):
    text = re.sub(r'\s+', ' ', text)  # Replace multiple spaces/newlines with a single space
    text = text.strip()  # Remove leading and trailing whitespace
    stop_words = set(stopwords.words('english'))
    word_tokens = word_tokenize(text)
    filtered_text = [word for word in word_tokens if word.lower() not in stop_words]
    text = ' '.join(filtered_text)
    return text

# Function to preprocess the image
def preprocess_image(image_path):
    input_image = Image.open(image_path).convert("RGB")
    preprocess = transforms.Compose([
        transforms.Resize(256),
        transforms.CenterCrop(224),
        transforms.ToTensor(),
        transforms.Normalize(mean=[0.485, 0.456, 0.406], std=[0.229, 0.224, 0.225]),
    ])
    input_tensor = preprocess(input_image)
    return input_tensor.unsqueeze(0)

# Function to extract and save images from PPT
def extract_images_from_ppt(ppt_path):
    prs = Presentation(ppt_path)
    image_files = []
    for slide_num, slide in enumerate(prs.slides):
        for shape_num, shape in enumerate(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = io.BytesIO(image.blob)
                image_file_name = f"slide_{slide_num+1}_image_{shape_num+1}.{image.ext}"
                with open(image_file_name, "wb") as f:
                    f.write(image_bytes.read())
                image_files.append(image_file_name)
    return image_files

# Function to extract text from an image using OCR
def extract_text_from_image(image_path):
    image = Image.open(image_path)
    return pytesseract.image_to_string(image)

# Extract text from PDF
def extract_text_from_pdf(pdf_path):
    with pdfplumber.open(pdf_path) as pdf:
        return "\n".join(page.extract_text() for page in pdf.pages)

# Extract text from PPT
def extract_text_from_ppt(ppt_path):
    prs = Presentation(ppt_path)
    return " ".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))

def read_word_file(file_path):
    doc = Document(file_path)
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)
    return '\n'.join(full_text)

# Extract limited images from PDF
def extract_limited_images_from_pdf(pdf_path, image_dir, limit=4):
    if not os.path.exists(image_dir):
        os.makedirs(image_dir)

    doc = fitz.open(pdf_path)
    image_captions = []
    image_count = 0

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        image_list = page.get_images(full=True)
        for img_index, img in enumerate(image_list):
            if image_count >= limit:
                return image_captions
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            image_path = os.path.join(image_dir, f"image{page_num}_{img_index}.png")
            with open(image_path, "wb") as f:
                f.write(image_bytes)
            image_captions.append(f"Image from page {page_num + 1}, image index {img_index + 1}")
            image_count += 1
    return image_captions

# Merge captions and text
def merge_captions_and_text(text, image_captions):
    return "\n\n".join([text] + image_captions)

def chunk_text(text, chunk_size=1000):
    """Splits the text into smaller chunks."""
    for i in range(0, len(text), chunk_size):
        yield text[i:i + chunk_size]

def summ(text, chunk_size=1000):
    summaries = []
    api_key = "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"
    client = groq.Client(api_key=api_key)

    for chunk in chunk_text(text, 1000):
        response = client.chat.completions.create(
            messages=[
                {
                    "role": "user",
                    "content": "Please summarize the following text:\n\n" + chunk,
                }
            ],
            model="llama3-8b-8192",
        )
        summaries.append(response.choices[0].message.content)

    # Combine the summaries of all chunks
    combined_summary = " ".join(summaries)
    return combined_summary

def format_bullet_points(text):
    formatted_response = text.replace('**', '')  # Add a newline before each numbered bullet point
    formatted_text = re.sub(r'\.\s*(\d+\.)', r'.\n\n\1', formatted_response)
    return formatted_text

# Answer question
def answer_question(question, context):
    api_key = "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"  # Replace with your Groq API key
    client = groq.Client(api_key=api_key)
    response = client.chat.completions.create(
        messages=[
            {
                "role": "user",
                "content": context + "\n\nQuestion: " + question,
            }
        ],
        model="llama3-8b-8192",
    )
    return response.choices[0].message.content

@app.route("/", methods=["GET", "POST"])
def home():
    status = "Leverage the chatbot's analysis tool for your documents.\nUpload your Document to access chat."
    result = None
    query = None
    result_sum = None
    history = QuestionAnswer.query.order_by(QuestionAnswer.timestamp.desc()).all()

    if request.method == "POST":
        if "file" in request.files:
            file = request.files["file"]
            if file.filename:
                # Save the uploaded file
                file_path = os.path.join("uploads", file.filename)
                file.save(file_path)

                # Determine the file type
                if file.filename.endswith(".pdf"):
                    text = extract_text_from_pdf(file_path)
                    summary = summ(text)
                    summary = format_bullet_points(summary)
                    new_summary = Summary(content=summary)
                    db.session.add(new_summary)
                    db.session.commit()

                    session['summary_id'] = new_summary.id
                    result_sum = summary
                    status = "PDF processed and summary generated."

                elif file.filename.endswith((".ppt", ".pptx")):
                    text = extract_text_from_ppt(file_path)
                    summary = summ(text)
                    extracted_images = extract_images_from_ppt(file_path)
                    descriptions = [summary]
                    for image_file in extracted_images:
                        extracted_text = extract_text_from_image(image_file)
                        descriptions.append(extracted_text)
                    result = "\n".join(descriptions)
                    summary = format_bullet_points(summary)
                    result_sum = Summary(content=summary)
                    db.session.add(result_sum)
                    db.session.commit()
                    session['summary_id'] = result_sum.id
                    status = "PPT processed and descriptions generated."

                elif file.filename.endswith((".doc", ".docx")):
                    text = read_word_file(file_path)
                    summary = summ(text)
                    new_summary = Summary(content=summary)
                    db.session.add(new_summary)
                    db.session.commit()

                    session['summary_id'] = new_summary.id
                    summary = format_bullet_points(summary)
                    result_sum = summary
                    status = "Word document processed and summary generated."

        elif "query" in request.form:
            query = request.form["query"]
            summary_id = session.get('summary_id')

            if summary_id:
                summary = Summary.query.get(summary_id)
                if summary:
                    context = summary.content
                    result = answer_question(query, context)
                    new_qa = QuestionAnswer(question=query, answer=result)
                    db.session.add(new_qa)
                    db.session.commit()
                    status = "Question answered based on the provided context."
                else:
                    status = "Summary not found for the provided context."
            else:
                status = "No summary available to provide context for the query."

    html = '''
    <!DOCTYPE html>
    <html lang="en">
    <head>
        <meta charset="UTF-8">
        <meta name="viewport" content="width=device-width, initial-scale=1.0">
        <title>Document Analyzer Chatbot</title>
        <link rel="stylesheet" href="https://stackpath.bootstrapcdn.com/bootstrap/4.5.2/css/bootstrap.min.css">
        <style>
            body {
                background-color: #f8f9fa;
                color: #333;
                font-family: Arial, sans-serif;
            }
            .container {
                margin-top: 30px;
                background-color: #fff;
                border-radius: 10px;
                padding: 30px;
                box-shadow: 0 0 10px rgba(0, 0, 0, 0.1);
            }
            .upload-form {
                margin-bottom: 30px;
            }
            .chat-window {
                max-height: 500px;
                overflow-y: auto;
                border: 1px solid #ccc;
                border-radius: 10px;
                padding: 10px;
                background-color: #e9ecef;
            }
            .chat-message {
                margin-bottom: 15px;
                padding: 10px;
                border-radius: 10px;
            }
            .user-query {
                background-color: #d4edda;
            }
            .user-response {
                background-color: #f8d7da;
            }
            .output-status {
                margin-top: 20px;
                font-style: italic;
                color: #666;
            }
            .upload-file{
                color: black;
                background-color:  #5a6268;
            }
            #query-input {
                width: 100%;
                padding: 10px;
                margin-top: 10px;
                border: 1px solid #ccc;
                border-radius: 10px;
            }
        </style>
    </head>
    <body>
        <div class="container">
            <h1 class="text-center mb-4">Document Analyzer Chatbot</h1>
            <form class="upload-form" method="POST" enctype="multipart/form-data">
                <div class="form-group">
                    <label for="file">Upload PDF, PPT, or Word Document:</label>
                    <input type="file" class="form-control-file" id="file" name="file">
                </div>
                <button type="submit" class="btn btn-primary btn-block upload-file">Upload and Process</button>
            </form>
            <form class="query-form" method="POST">
                <div class="form-group">
                    <label for="query">Ask a Question:</label>
                    <input type="text" class="form-control" id="query-input" name="query">
                </div>
                <button type="submit" class="btn btn-success btn-block">Submit Question</button>
            </form>
            <div class="chat-window" id="chat-window">
                <p id="status-message"></p>
                {% if query %}
                <div class="chat-message user-query" id="boxing">
                    <p style="background-color:rgb(4,4,102);color:#ffff;border-radius:20px 20px 20px 20px;padding:15px"><strong>Question: {{ query }}</strong> </p>
                </div>
                {% endif %}
                {% if result_sum %}
                    <p style="background-color:rgb(4,4,102);color:#ffff;border-radius:20px 20px 20px 20px;padding:15px"><b>{{ result_sum }}</b></p><br><br>
                {% endif %}
                {% if result %}
                <div class="chat-message user-response">
                    <p style="background-color:rgb(4,4,102);color:#ffff;border-radius:20px 20px 20px 20px;padding:15px"><strong>Response-> {{ result }}</strong> </p>
                </div>
                {% endif %}
                <div class="output-status">Upload the PDF, PPT, or Word and get its analysis and ask your questions{{ status }}</div>
            </div>
            <h3 class="mt-4">History</h3>
            <ul class="list-group">
                {% for qa in history %}
                <li class="list-group-item">
                    <strong>Question:</strong> {{ qa.question }}<br>
                    <strong>Answer:</strong> {{ qa.answer }}
                </li>
                {% endfor %}
            </ul>
        </div>
    </body>
    </html>
    '''
    return render_template_string(html, result=result, query=query, result_sum=result_sum, status=status, history=history)

if __name__ == '__main__':
    # Run the Flask application
    app.run(host='0.0.0.0', port=int(os.environ.get('PORT', 5000)))
